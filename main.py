import pandas as pd
import re
import os
import json
from PIL import Image
from logging import PlaceHolder
from pptx import Presentation
from pptx.util import Pt

with open("config.json") as json_file:
    data = json.load(json_file)

    DisplayCol = data["column"]['displayName']
    YearCol = data["column"]['year']
    MajorCol = data["column"]['major']
    CaptionCol = data["column"]['caption']

    ExcelLocation = data["location"]["excel"]
    TemplateLocation = data["location"]["template"]
    FontLocation = data["location"]["font"]
    PhotoLocation = data["location"]["photo"]
    PptxDestination = data["location"]["target"]


def ProcessField(s):
    return re.sub(r'(?<!\w)and(?!\w)', "&", str(s)).upper()


def create_name_mapping(df, display_col):
    """
    Create a mapping between Excel display names and image filenames to handle edge cases.
    This helps with cases where display names might be stored differently in Excel vs image filenames.
    """
    name_mapping = {}
    picArr = [f for f in os.listdir(PhotoLocation) 
              if f.endswith(('.jpg', '.jpeg', '.png')) and not f.startswith('.')]
    
    print("Creating display name mapping...")
    for index, row in df.iterrows():
        display_name = str(row[display_col]).strip()
        if pd.isna(display_name) or display_name == '':
            continue
            
        # Try to find the best match for this display name
        best_match = GetFileName(display_name)
        if best_match:
            name_mapping[display_name] = best_match
            print(f"Mapped '{display_name}' -> '{best_match}'")
        else:
            print(f"WARNING: No image found for '{display_name}'")
    
    return name_mapping


def GetFileName(name):
    """
    Robust name matching function that handles edge cases like multiple people with the same first name.
    Uses a scoring system to find the best match.
    """
    name = name.strip()
    picArr = [f for f in os.listdir(PhotoLocation) 
              if f.endswith(('.jpg', '.jpeg', '.png')) and not f.startswith('.')]
    
    if not picArr:
        print(f"No image files found in {PhotoLocation}")
        return None
    
    # Clean the name for matching
    clean_name = re.sub(r'[^\w\s]', '', name.lower()).strip()
    name_words = [word for word in clean_name.split() if len(word) > 1]
    
    if not name_words:
        print(f"Invalid name format: {name}")
        return None
    
    best_match = None
    best_score = 0
    matches = []
    
    # Check if this is a unique first name (helpful for first-name-only cases)
    first_name = name_words[0]
    first_name_count = sum(1 for f in picArr 
                          if re.sub(r'[^\w\s]', '', f.lower()).strip().startswith(first_name))
    
    for filename in picArr:
        # Clean the filename for comparison
        clean_filename = re.sub(r'[^\w\s]', '', filename.lower()).strip()
        clean_filename_no_ext = os.path.splitext(clean_filename)[0]
        filename_words = [word for word in clean_filename_no_ext.split() if len(word) > 1]
        
        if not filename_words:
            continue
        
        # Calculate match score using multiple criteria
        score = 0
        
        # 1. Exact match gets highest score
        if clean_name == clean_filename_no_ext:
            score = 1000
        # 2. Exact match without extension
        elif clean_name == clean_filename_no_ext:
            score = 900
        # 3. All name words found in filename (in order) - this is crucial for multiple people with same first name
        elif all(word in clean_filename_no_ext for word in name_words):
            score = 800 - len(name_words)  # Prefer shorter names for same score
            
            # Bonus for exact word order match
            if ' '.join(name_words) in clean_filename_no_ext:
                score += 200
        # 4. Handle cases where only first name is in filename
        elif len(name_words) > 1 and name_words[0] in clean_filename_no_ext:
            # If first name matches and it's a single word filename, this might be a first-name-only case
            if len(filename_words) == 1 and name_words[0] == filename_words[0]:
                if first_name_count == 1:
                    score = 700  # High score for exact first name match when it's unique
                else:
                    score = 300  # Lower score when multiple people have same first name
            elif name_words[0] in clean_filename_no_ext:
                # Check if this is the only person with this first name
                if first_name_count == 1:
                    score = 600  # Only person with this first name, so it's likely correct
                else:
                    # Multiple people with same first name, need more specific matching
                    score = 200
        # 5. Most name words found in filename
        else:
            matching_words = sum(1 for word in name_words if word in clean_filename_no_ext)
            if matching_words > 0:
                score = matching_words * 100
                
                # Bonus for consecutive word matches
                consecutive_bonus = 0
                for i in range(len(name_words) - 1):
                    if (name_words[i] in clean_filename_no_ext and 
                        name_words[i + 1] in clean_filename_no_ext):
                        consecutive_bonus += 50
                score += consecutive_bonus
                
                # Additional penalty for partial matches when we have multiple people with same first name
                if len(name_words) > 1 and matching_words < len(name_words):
                    score -= 100  # Heavy penalty for incomplete matches when we have multiple words
        
        # Penalty for extra words in filename (to avoid matching "Chloe" to "Chloe Ng" when looking for "Chloe Siew")
        extra_words = len(filename_words) - len(name_words)
        if extra_words > 0:
            score -= extra_words * 20
        
        # Store match if score is above threshold
        if score > 0:
            matches.append((filename, score))
            if score > best_score:
                best_score = score
                best_match = filename
    
    if not matches:
        print(f"No picture found for '{name}' (cleaned: '{clean_name}')")
        return None
    
    # If we have multiple matches with the same score, we need to be more specific
    if len(matches) > 1:
        max_score = max(score for _, score in matches)
        high_score_matches = [f for f, s in matches if s == max_score]
        
        if len(high_score_matches) > 1:
            print(f"Multiple high-confidence matches found for '{name}':")
            for match in high_score_matches:
                print(f"  - {match}")
            print(f"Choosing first match: {high_score_matches[0]}")
            return high_score_matches[0]
    
    return best_match


# def PrimePics():
#     imagenames = os.listdir(PhotoLocation)
#     for imagename in imagenames:
#         path = os.path.join(PhotoLocation, imagename)
#         try:
#             im = Image.open(path)
#             im.save(path)
#         except Exception as e:
#             print(e)

def PrimePics():
    imagenames = os.listdir(PhotoLocation)
    for imagename in imagenames:
        if imagename.startswith('.'):  # Skip hidden files like .DS_Store
            continue
        if not (imagename.endswith('.jpg') or imagename.endswith('.jpeg') or imagename.endswith('.png')):
            continue
            
        path = os.path.join(PhotoLocation, imagename)
        try:
            with Image.open(path) as im:
                # Convert to RGB if necessary
                if im.mode in ['RGBA', 'LA', 'P']:
                    im = im.convert('RGB')
                # Save as JPEG with better quality
                im.save(path, format='JPEG', quality=95, optimize=True)
        except Exception as e:
            print(f"Error processing image {imagename}: {e}")
            # Try to remove the problematic file
            try:
                os.remove(path)
                print(f"Removed problematic file: {imagename}")
            except:
                pass


def count_existing_pptx_files(df, display_col):
    """
    Count how many PPTX files already exist for the given display names.
    """
    existing_count = 0
    total_count = 0
    
    for index, row in df.iterrows():
        display_name = str(row[display_col]).strip()
        if pd.isna(display_name) or display_name == '':
            continue
        
        total_count += 1
        if check_pptx_exists(display_name, False):  # Always check without force_recreate for counting
            existing_count += 1
    
    return existing_count, total_count


def check_pptx_exists(name, force_recreate=False):
    """
    Check if a PPTX file already exists for the given name.
    If force_recreate is True, always return False (file doesn't exist).
    """
    if force_recreate:
        return False
    
    safeName = re.sub(r'[^A-z]', "", str(name)) + "_Noctua.pptx"
    pptx_path = os.path.join(PptxDestination, safeName)
    return os.path.exists(pptx_path)


def CreateDoorcard(name, data_dict, name_mapping=None, force_recreate=False):
    # Check if PPTX already exists
    if check_pptx_exists(name, force_recreate):
        return "skipped"
    
    prs = Presentation(TemplateLocation)
    phs = prs.slides[0].placeholders
    
    for ph in phs:
        if ph.name == "Picture":
            # Use mapping if available, otherwise use GetFileName
            if name_mapping and name in name_mapping:
                filename = name_mapping[name]
            else:
                filename = GetFileName(name)
            
            if filename is None:
                print(f"Skipping {name} - no image found")
                return False
            try:
                ph.insert_picture(os.path.join(PhotoLocation, filename))
            except Exception as e:
                print(f"Error inserting picture for {name}: {e}")
                return False
        else:
            try:
                ph.text = ProcessField(data_dict[ph.name])
            except Exception as e:
                print(f"Error processing field {ph.name} for {name}: {e}")
                ph.text = str(data_dict.get(ph.name, ""))

    safeName = re.sub(r'[^A-z]', "", str(name)) + "_Noctua.pptx"
    if not os.path.exists(PptxDestination):
        os.makedirs(PptxDestination)
    prs.save(os.path.join(PptxDestination, safeName))
    return True


def validate_mapping(name_mapping, df, display_col):
    """
    Validate the name mapping and identify potential issues.
    """
    print("\nValidating display name mapping...")
    
    # Check for duplicate mappings (same image used for multiple people)
    image_to_names = {}
    for display_name, image_name in name_mapping.items():
        if image_name in image_to_names:
            image_to_names[image_name].append(display_name)
        else:
            image_to_names[image_name] = [display_name]
    
    # Report duplicates
    duplicates = {img: names for img, names in image_to_names.items() if len(names) > 1}
    if duplicates:
        print("WARNING: Multiple people mapped to the same image:")
        for img, names in duplicates.items():
            print(f"  Image: {img}")
            for name in names:
                print(f"    - {name}")
        
        print("\nThese duplicates likely occur because:")
        print("1. Multiple people have the same first name and only used their first name in the image filename")
        print("2. The matching algorithm couldn't distinguish between similar names")
        print("3. There might be naming inconsistencies in the Excel file vs image filenames")
    
    # Check for unmapped names
    all_names = set(str(name).strip() for name in df[display_col].dropna())
    mapped_names = set(name_mapping.keys())
    unmapped = all_names - mapped_names
    
    if unmapped:
        print(f"\nWARNING: {len(unmapped)} display names could not be mapped to images:")
        
        # Group unmapped names by first name to see patterns
        first_name_groups = {}
        for name in unmapped:
            first_name = name.split()[0].lower() if name.split() else name.lower()
            if first_name not in first_name_groups:
                first_name_groups[first_name] = []
            first_name_groups[first_name].append(name)
        
        # Show grouped unmapped names
        for first_name, names in sorted(first_name_groups.items()):
            if len(names) > 1:
                print(f"  Multiple people with first name '{first_name.title()}':")
                for name in sorted(names):
                    print(f"    - {name}")
            else:
                print(f"  - {names[0]}")
    
    return len(duplicates) == 0 and len(unmapped) == 0


def handle_first_name_only_cases(name_mapping, df, display_col):
    """
    Handle cases where people only put their first name in the image filename.
    This function tries to resolve conflicts by looking for more specific matches.
    """
    print("\nHandling first-name-only cases for display names...")
    
    # Get all images that might be first-name-only
    picArr = [f for f in os.listdir(PhotoLocation) 
              if f.endswith(('.jpg', '.jpeg', '.png')) and not f.startswith('.')]
    
    # Group names by first name
    first_name_groups = {}
    for index, row in df.iterrows():
        display_name = str(row[display_col]).strip()
        if pd.isna(display_name) or display_name == '':
            continue
        
        first_name = display_name.split()[0].lower() if display_name.split() else display_name.lower()
        if first_name not in first_name_groups:
            first_name_groups[first_name] = []
        first_name_groups[first_name].append(display_name)
    
    # Look for cases where multiple people have the same first name
    conflicts = {first_name: names for first_name, names in first_name_groups.items() if len(names) > 1}
    
    if conflicts:
        print("Found multiple people with the same first name:")
        for first_name, names in conflicts.items():
            print(f"  {first_name.title()}: {', '.join(names)}")
            
            # Look for images that might match these people more specifically
            matching_images = []
            for img in picArr:
                clean_img = re.sub(r'[^\w\s]', '', img.lower()).strip()
                clean_img_no_ext = os.path.splitext(clean_img)[0]
                
                # Check if this image matches any of the names more specifically
                for name in names:
                    clean_name = re.sub(r'[^\w\s]', '', name.lower()).strip()
                    if clean_name in clean_img_no_ext or any(word in clean_img_no_ext for word in clean_name.split()):
                        matching_images.append((img, name))
            
            if matching_images:
                print(f"    Potential matches found:")
                for img, name in matching_images:
                    print(f"      {name} -> {img}")
    
    return name_mapping


if __name__ == "__main__":
    df = pd.read_excel(ExcelLocation)
    PrimePics()
    
    # Count existing PPTX files
    existing_count, total_count = count_existing_pptx_files(df, DisplayCol)
    print(f"\nFound {existing_count} existing PPTX files out of {total_count} total entries")
    print(f"Will create {total_count - existing_count} new PPTX files")
    
    # Ask user if they want to force recreate all files
    if existing_count > 0:
        print(f"\nSome PPTX files already exist. Options:")
        print("1. Skip existing files (recommended)")
        print("2. Force recreate all files")
        response = input("Choose option (1 or 2): ").strip()
        force_recreate = response == "2"
        if force_recreate:
            print("Will force recreate all files (existing files will be overwritten)")
        else:
            print("Will skip existing files and only create new ones")
    else:
        force_recreate = False
        print("No existing PPTX files found. Will create all files.")
    
    # Create name mapping to handle edge cases
    name_mapping = create_name_mapping(df, DisplayCol)
    
    # Handle first-name-only cases
    name_mapping = handle_first_name_only_cases(name_mapping, df, DisplayCol)
    
    # Validate the mapping
    mapping_valid = validate_mapping(name_mapping, df, DisplayCol)
    
    if not mapping_valid:
        print("\nWARNING: Issues found in name mapping. Please review the warnings above.")
        print("\nTo resolve these issues, you may need to:")
        print("1. Check if image filenames use full names vs first names only")
        print("2. Manually rename some image files to include full names")
        print("3. Update the Excel file to match image naming conventions")
        print("4. Or continue and manually fix any incorrect assignments later")
        
        response = input("\nContinue anyway? (y/n): ")
        if response.lower() != 'y':
            print("Exiting...")
            exit(1)
    
    total = len(df)
    success = 0
    skipped = 0
    
    print(f"\nProcessing {total} doorcards...")
    
    for i, row in df.iterrows():
        try:
            display_name = str(row[DisplayCol]).strip()
            if pd.isna(display_name) or display_name == '':
                print(f"Skipping row {i}: Empty display name")
                continue
            
            # Use the mapping if available, otherwise fall back to GetFileName
            if display_name in name_mapping:
                image_filename = name_mapping[display_name]
            else:
                image_filename = GetFileName(display_name)
                if image_filename:
                    name_mapping[display_name] = image_filename
            
            if not image_filename:
                print(f"Skipping {display_name} - no image found")
                continue
            
            result = CreateDoorcard(
                display_name,
                {
                    "Name": row[DisplayCol],
                    "Year": row[YearCol],
                    "Major": row[MajorCol],
                    "Caption": row[CaptionCol]
                },
                name_mapping,
                force_recreate
            )
            
            if result == "skipped":
                skipped += 1
                print(f"Skipped {display_name} - PPTX already exists")
            elif result == True:
                success += 1
                print(f"Created doorcard for {display_name} using image: {image_filename}")
            else:
                print(f"Failed to create doorcard for {display_name}")
                
        except Exception as e:
            print(f"Error at {row[DisplayCol] if 'row' in locals() and DisplayCol in row else 'unknown'}: {e}")
    
    print(f"\nCreation completed: {success} new / {skipped} skipped / {total} total")
    
    # Print summary of unmapped names
    unmapped = [name for name in df[DisplayCol].dropna() if str(name).strip() not in name_mapping]
    if unmapped:
        print(f"\nWARNING: {len(unmapped)} display names could not be mapped to images:")
        for name in unmapped:
            print(f"  - {name}")
